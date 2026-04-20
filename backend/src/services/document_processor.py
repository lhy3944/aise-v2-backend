"""문서 파싱 + 처리 파이프라인"""

import uuid

from loguru import logger
from sqlalchemy import select

from src.core.database import async_session
from src.core.exceptions import AppException
from src.models.knowledge import KnowledgeDocument, KnowledgeChunk
from src.services import storage_svc, embedding_svc
from src.utils.text_chunker import _get_encoding, chunk_text


def parse_document(file_bytes: bytes, file_type: str) -> str:
    """파일 바이트를 텍스트로 변환한다."""
    file_type = file_type.lower()

    if file_type in ("txt", "md"):
        return file_bytes.decode("utf-8", errors="replace")

    if file_type == "pdf":
        return _parse_pdf(file_bytes)

    if file_type == "docx":
        return _parse_docx(file_bytes)

    if file_type == "pptx":
        return _parse_pptx(file_bytes)

    if file_type == "xlsx":
        return _parse_xlsx(file_bytes)

    raise AppException(400, f"지원하지 않는 파일 형식입니다: {file_type}")


def _parse_pdf(data: bytes) -> str:
    import pymupdf

    text_parts = []
    with pymupdf.open(stream=data, filetype="pdf") as doc:
        for page in doc:
            text_parts.append(page.get_text())
    return "\n".join(text_parts)


def _parse_docx(data: bytes) -> str:
    import io
    from docx import Document

    doc = Document(io.BytesIO(data))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _parse_pptx(data: bytes) -> str:
    import io
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    text_parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        text_parts.append(text)
    return "\n".join(text_parts)


def _parse_xlsx(data: bytes) -> str:
    import io
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    text_parts = []
    for ws in wb.worksheets:
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                text_parts.append("\t".join(cells))
    wb.close()
    return "\n".join(text_parts)


async def process_document(document_id: uuid.UUID) -> None:
    """문서 처리 파이프라인: 다운로드 -> 파싱 -> 청킹 -> 임베딩 -> DB 저장

    BackgroundTasks에서 호출되며, 독립된 DB 세션을 생성하여 사용한다.
    에러 시 document status를 'failed'로 업데이트한다.
    """
    logger.info(f"문서 처리 시작: document_id={document_id}")

    async with async_session() as db:
        # 1. DB에서 문서 조회
        result = await db.execute(
            select(KnowledgeDocument).where(KnowledgeDocument.id == document_id)
        )
        doc = result.scalar_one_or_none()
        if not doc:
            logger.error(f"문서를 찾을 수 없음: document_id={document_id}")
            return

        try:
            # 2. MinIO에서 파일 다운로드
            bucket = storage_svc.get_default_bucket()
            file_bytes = await storage_svc.download_file(bucket, doc.storage_key)

            # 3. 파싱
            text = parse_document(file_bytes, doc.file_type)
            if not text.strip():
                raise ValueError("문서에서 텍스트를 추출할 수 없습니다.")

            # 4. 청킹
            chunks = chunk_text(text, max_tokens=500, overlap_tokens=50, file_type=doc.file_type)
            if not chunks:
                raise ValueError("청킹 결과가 비어 있습니다.")

            # 5. 임베딩
            embeddings = await embedding_svc.get_embeddings(chunks)

            # 6. KnowledgeChunk 레코드 생성
            encoding = _get_encoding()

            for i, (chunk_text_content, embedding) in enumerate(zip(chunks, embeddings)):
                chunk = KnowledgeChunk(
                    document_id=doc.id,
                    project_id=doc.project_id,
                    chunk_index=i,
                    content=chunk_text_content,
                    token_count=len(encoding.encode(chunk_text_content)),
                    embedding=embedding,
                    metadata_={"document_name": doc.name},
                )
                db.add(chunk)

            # 7. 문서 상태 업데이트
            doc.status = "completed"
            doc.chunk_count = len(chunks)
            await db.commit()

            logger.info(f"문서 처리 완료: document_id={document_id}, chunks={len(chunks)}")

        except Exception as e:
            logger.error(f"문서 처리 실패: document_id={document_id}, error={e}")
            await db.rollback()

            # 에러 상태로 업데이트 (새 트랜잭션)
            result = await db.execute(
                select(KnowledgeDocument).where(KnowledgeDocument.id == document_id)
            )
            doc = result.scalar_one_or_none()
            if doc:
                doc.status = "failed"
                doc.error_message = str(e)[:500]
                await db.commit()
